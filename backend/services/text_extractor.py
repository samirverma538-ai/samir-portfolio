from pathlib import Path


def extract_text(file_path: Path, file_type: str) -> str:
    ext = file_type.lower()
    try:
        if ext == ".pdf":
            return _extract_pdf(file_path)
        if ext in (".doc", ".docx"):
            return _extract_docx(file_path)
        if ext in (".ppt", ".pptx"):
            return _extract_pptx(file_path)
        if ext in (".jpg", ".jpeg", ".png"):
            return f"[Image file: {file_path.name}]"
        if ext == ".txt":
            return file_path.read_text(encoding="utf-8", errors="replace")
    except Exception as exc:
        return f"[Text extraction failed: {exc}]"
    return ""


def _extract_pdf(path: Path) -> str:
    from PyPDF2 import PdfReader

    reader = PdfReader(str(path))
    parts = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            parts.append(text)
    return "\n".join(parts)


def _extract_docx(path: Path) -> str:
    from docx import Document

    doc = Document(str(path))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        if slide_text:
            parts.append(f"Slide {i}: " + " | ".join(slide_text))
    return "\n".join(parts)
